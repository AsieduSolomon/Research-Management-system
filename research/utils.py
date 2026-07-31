import re
import math
import random
from collections import Counter
from django.conf import settings
import os

# Import file processing libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

STOP_WORDS = {
    'a', 'an', 'and', 'are', 'as', 'at', 'be', 'by', 'for', 'from', 'has', 'he',
    'in', 'is', 'it', 'its', 'of', 'on', 'that', 'the', 'to', 'was', 'were',
    'will', 'with', 'this', 'that', 'these', 'those', 'but', 'or', 'so', 'for',
    'can', 'may', 'could', 'should', 'would', 'have', 'been', 'being', 'am',
    'not', 'no', 'yes', 'why', 'how', 'what', 'when', 'where', 'who', 'which',
    'research', 'study', 'paper', 'proposal', 'also', 'however', 'therefore'
}


def extract_text_from_pptx(file_path):
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        return " ".join(text_parts)
    except Exception:
        return ""


def extract_text_from_docx(file_path):
    if not DOCX_AVAILABLE:
        return ""
    try:
        doc = Document(file_path)
        text_parts = [p.text for p in doc.paragraphs if p.text]
        return " ".join(text_parts)
    except Exception:
        return ""


def extract_text_from_pdf(file_path):
    if not PDF_AVAILABLE:
        return ""
    try:
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return " ".join(text_parts)
    except Exception:
        return ""


def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        return ""


def extract_text_from_file(file_path, file_extension):
    ext = file_extension.lower()
    if ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    return ""


def generate_keywords_from_text(text, max_keywords=15):
    if not text or len(text.strip()) < 10:
        return "research proposal"

    words = re.findall(r'\b[a-z]{3,}\b', text.lower())
    filtered_words = [w for w in words if w not in STOP_WORDS and len(w) > 2]

    if not filtered_words:
        return "research proposal"

    word_counts = Counter(filtered_words)
    top_words = [word for word, count in word_counts.most_common(max_keywords)]

    return ", ".join(top_words) if top_words else "research proposal"


def _tokenize(text):
    """Tokenise text into a Counter of lowercased alphabetic tokens ≥3 chars,
    excluding stop-words."""
    if not text:
        return Counter()
    tokens = re.findall(r'\b[a-z]{3,}\b', text.lower())
    return Counter(t for t in tokens if t not in STOP_WORDS)


def _tfidf_vector(term_counts, idf):
    """Return a TF-IDF dict given a raw term Counter and a pre-computed IDF dict."""
    total = sum(term_counts.values()) or 1
    return {term: (count / total) * idf.get(term, 1.0)
            for term, count in term_counts.items()}


def _cosine_similarity(vec_a, vec_b):
    """Cosine similarity between two TF-IDF dicts."""
    if not vec_a or not vec_b:
        return 0.0
    dot = sum(vec_a.get(t, 0) * vec_b.get(t, 0) for t in vec_a)
    norm_a = math.sqrt(sum(v * v for v in vec_a.values()))
    norm_b = math.sqrt(sum(v * v for v in vec_b.values()))
    denom = norm_a * norm_b
    return dot / denom if denom > 0 else 0.0


def compute_interest_match_matrix(student_ids, supervisor_ids):
    """
    Build a (student_id, supervisor_id) → cosine-similarity score dict.

    Uses TF-IDF weighting over the corpus of student proposals + supervisor
    specialisations so that rare domain-specific terms carry more weight than
    common words.
    """
    from .models import Proposal, User

    # ----- gather raw text -----
    student_texts = {}
    for sid in student_ids:
        try:
            proposal = (Proposal.objects
                        .filter(student_id=sid, status='approved')
                        .order_by('-created_at')
                        .first())
            if proposal:
                student_texts[sid] = (
                    f"{proposal.keywords or ''} "
                    f"{proposal.title or ''} "
                    f"{proposal.abstract or ''} "
                    f"{proposal.objectives or ''}"
                )
            else:
                student_texts[sid] = ""
        except Exception:
            student_texts[sid] = ""

    supervisor_texts = {}
    for sup_id in supervisor_ids:
        try:
            sup = User.objects.get(id=sup_id)
            parts = [sup.department or ""]
            if hasattr(sup, 'supervisor_profile'):
                parts.append(sup.supervisor_profile.specialization or "")
                parts.append(sup.supervisor_profile.bio or "")
                parts.append(sup.supervisor_profile.qualifications or "")
            supervisor_texts[sup_id] = " ".join(parts)
        except Exception:
            supervisor_texts[sup_id] = ""

    # ----- build IDF over entire corpus -----
    all_counters = []
    student_counters = {sid: _tokenize(txt) for sid, txt in student_texts.items()}
    supervisor_counters = {sup_id: _tokenize(txt) for sup_id, txt in supervisor_texts.items()}
    all_counters = list(student_counters.values()) + list(supervisor_counters.values())

    # document frequency
    df = Counter()
    for c in all_counters:
        for term in c:
            df[term] += 1
    n_docs = max(len(all_counters), 1)
    idf = {term: math.log((n_docs + 1) / (freq + 1)) + 1.0
           for term, freq in df.items()}

    # ----- compute TF-IDF vectors & cosine similarity -----
    match_matrix = {}
    for sid in student_ids:
        vec_s = _tfidf_vector(student_counters[sid], idf)
        for sup_id in supervisor_ids:
            vec_sup = _tfidf_vector(supervisor_counters[sup_id], idf)
            score = _cosine_similarity(vec_s, vec_sup)
            match_matrix[(sid, sup_id)] = score

    return match_matrix


def _logistic(x, r=3.9):
    """Single step of the logistic map."""
    return r * x * (1.0 - x)


def get_chaotic_allocation(proposal_ids, student_ids, supervisor_ids,
                           supervisor_capacities, interest_match_matrix,
                           w1=0.6, w2=0.4, max_iterations=200):
    """
    Chaotic Logistic-Map Allocation Algorithm
    ==========================================
    Uses the logistic map  x_{n+1} = r·x_n·(1 - x_n)  (r = 3.9, chaotic regime)
    to generate diverse, pseudo-random orderings of students each iteration.
    Each iteration produces a candidate allocation; the best-fitness allocation
    across all iterations is returned.

    Fitness (per candidate):
        F = w1 * avg_match_score  +  w2 * load_balance_score

    where
        avg_match_score    = mean cosine similarity (student ↔ supervisor)
        load_balance_score = 1 - (max_load - min_load) / max(max_load, 1)

    Parameters
    ----------
    proposal_ids         : list[int]  – proposal pk list (parallel to student_ids)
    student_ids          : list[int]  – student pk list
    supervisor_ids       : list[int]
    supervisor_capacities: dict[sup_id → remaining capacity]
    interest_match_matrix: dict[(student_id, sup_id) → float 0-1]
    w1, w2               : float weights (must sum to 1)
    max_iterations       : int

    Returns
    -------
    list[(proposal_id, student_id, supervisor_id, match_score)]
    """
    from .models import Allocation

    if not student_ids or not supervisor_ids:
        return []

    # Base workload (already-active allocations)
    base_workload = {}
    for sup_id in supervisor_ids:
        base_workload[sup_id] = (
            Allocation.objects.filter(supervisor_id=sup_id, status='active').count()
        )

    n_students = len(student_ids)
    indices = list(range(n_students))

    # Seed multiple chaotic trajectories from different starting points
    # to improve diversity across iterations.
    seeds = [0.1 + 0.05 * i for i in range(5)]  # 5 independent chaotic seeds
    seed_idx = 0

    r = 3.9
    x_values = list(seeds)  # maintain independent trajectories

    best_allocation = None
    best_fitness = -1.0

    for iteration in range(max_iterations):
        # Advance one chaotic trajectory per iteration (round-robin)
        traj = seed_idx % len(x_values)
        x_values[traj] = _logistic(x_values[traj], r)
        chaos_value = x_values[traj]
        seed_idx += 1

        # Use chaos value as an additional shuffle seed so each iteration
        # produces a genuinely different ordering.
        rng = random.Random(int(chaos_value * 1_000_000) ^ (iteration * 2654435761))
        shuffled = indices[:]
        rng.shuffle(shuffled)

        candidate = []
        workload = base_workload.copy()
        remaining_cap = {s: supervisor_capacities.get(s, 5) for s in supervisor_ids}

        for idx in shuffled:
            prop_id = proposal_ids[idx]
            s_id = student_ids[idx]
            best_sup = None
            best_score = -1.0

            # Score every supervisor and pick the best feasible one
            for sup_id in supervisor_ids:
                if remaining_cap.get(sup_id, 0) <= 0:
                    continue

                match_score = interest_match_matrix.get((s_id, sup_id), 0.0)

                # Load balance term: penalise giving more students to an already
                # busier supervisor relative to ideal even distribution.
                total_capacity = sum(supervisor_capacities.get(s, 5)
                                     for s in supervisor_ids)
                ideal_load = n_students / len(supervisor_ids)
                projected_load = workload.get(sup_id, 0) + 1
                max_cap = supervisor_capacities.get(sup_id, 5)
                load_ratio = projected_load / max(max_cap, 1)
                # balance_score = 1 when load_ratio = 0, 0 when fully loaded
                balance_score = 1.0 - load_ratio

                fitness = w1 * match_score + w2 * balance_score

                if fitness > best_score:
                    best_score = fitness
                    best_sup = sup_id

            if best_sup is not None:
                candidate.append((prop_id, s_id, best_sup, best_score))
                workload[best_sup] = workload.get(best_sup, 0) + 1
                remaining_cap[best_sup] -= 1

        if not candidate:
            continue

        # Evaluate candidate fitness
        avg_match = sum(score for _, _, _, score in candidate) / len(candidate)
        loads = list(workload.values())
        max_load = max(loads)
        min_load = min(loads)
        balance_score = 1.0 - ((max_load - min_load) / max(max_load, 1)) if max_load > 0 else 1.0
        overall_fitness = w1 * avg_match + w2 * balance_score

        if overall_fitness > best_fitness:
            best_fitness = overall_fitness
            best_allocation = candidate

    return best_allocation or []
